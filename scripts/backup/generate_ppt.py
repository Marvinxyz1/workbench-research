# -*- coding: utf-8 -*-
"""
KPMG Workbench 深度概要 PPT 生成器
基于 HTML 幻灯片内容生成 PowerPoint 文件
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# 设置 UTF-8 编码输出
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding='utf-8')

# KPMG 品牌颜色常量
KPMG_BLUE = RGBColor(0, 51, 141)  # #00338D
KPMG_DARK_BLUE = RGBColor(0, 51, 141)
LIGHT_GRAY = RGBColor(240, 244, 248)
BRIGHT_BLUE = RGBColor(0, 123, 255)
TEXT_PRIMARY = RGBColor(26, 26, 26)
TEXT_SECONDARY = RGBColor(85, 85, 85)
BORDER_COLOR = RGBColor(221, 228, 233)
WHITE = RGBColor(255, 255, 255)
RED_BG = RGBColor(254, 242, 242)
RED_BORDER = RGBColor(254, 202, 202)
RED_TEXT = RGBColor(185, 28, 28)
GREEN_BG = RGBColor(240, 253, 244)
GREEN_BORDER = RGBColor(187, 247, 208)
GREEN_TEXT = RGBColor(21, 128, 61)
BLUE_BG = RGBColor(239, 246, 255)


def set_font(text_frame, font_name="微软雅黑", font_size=18, bold=False, color=TEXT_SECONDARY):
    """设置文本框的字体样式"""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=TEXT_SECONDARY, alignment=PP_ALIGN.LEFT):
    """添加文本框"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True

    for paragraph in text_frame.paragraphs:
        paragraph.alignment = alignment
        for run in paragraph.runs:
            run.font.name = "微软雅黑"
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color

    return textbox


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=TEXT_SECONDARY, bullet_color=KPMG_BLUE):
    """添加项目符号列表"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(items):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_before = Pt(6) if i > 0 else Pt(0)

        for run in p.runs:
            run.font.name = "微软雅黑"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color

    return textbox


def add_card(slide, left, top, width, height, bg_color=LIGHT_GRAY, border_color=BORDER_COLOR):
    """添加卡片背景"""
    card = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1)
    return card


def add_footer(slide, left_text="KPMG Workbench 深度概要", right_text="", y_position=6.8):
    """添加页脚"""
    # 左侧文本
    add_text_box(slide, 0.5, y_position, 4, 0.3, left_text, font_size=10, color=RGBColor(136, 136, 136))
    # 右侧文本
    if right_text:
        add_text_box(slide, 6, y_position, 4, 0.3, right_text, font_size=10, color=RGBColor(136, 136, 136), alignment=PP_ALIGN.RIGHT)


def create_slide_1(prs):
    """幻灯片 1: 封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 装饰线
    line = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(1.8), Inches(1.04), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = KPMG_BLUE
    line.line.fill.background()

    # 主标题
    add_text_box(slide, 0.5, 2.0, 9, 0.8, "KPMG Workbench 深度概要",
                 font_size=44, bold=True, color=KPMG_BLUE)

    # 副标题
    add_text_box(slide, 0.5, 2.9, 8, 0.5, "解构 KPMG 的 \"AI 骨干系统\" (AI Backbone)",
                 font_size=24, bold=False, color=TEXT_SECONDARY)

    # 描述文本
    desc = "我们将从三个核心维度（战略、技术、生态）深入剖析 KPMG Workbench，揭示这一平台如何重塑 KPMG 的 AI 创新能力。"
    add_text_box(slide, 0.5, 4.2, 7, 1, desc, font_size=18, color=TEXT_SECONDARY)

    # 页脚
    add_text_box(slide, 0.5, 6.8, 5, 0.3, "调研日期: 2025-11-17 | 调研人员: Claude Code",
                 font_size=10, color=RGBColor(136, 136, 136))
    add_text_box(slide, 5, 6.8, 4.5, 0.3, "信息来源: KPMG Workbench 官方 SharePoint 站点",
                 font_size=10, color=RGBColor(136, 136, 136), alignment=PP_ALIGN.RIGHT)


def create_slide_2(prs):
    """幻灯片 2: 执行摘要"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    add_text_box(slide, 0.5, 0.5, 9, 0.6, "执行摘要", font_size=44, bold=True, color=KPMG_BLUE)
    add_text_box(slide, 0.5, 1.1, 9, 0.4, "Workbench 是 KPMG 的核心 AI 战略资产",
                 font_size=24, color=TEXT_SECONDARY)

    # 核心洞察卡片
    add_card(slide, 0.5, 1.8, 9, 1.3, bg_color=BLUE_BG, border_color=BRIGHT_BLUE)
    add_text_box(slide, 0.6, 1.95, 8.8, 0.3, "核心洞察：从\"碎片化\"到\"一体化\"",
                 font_size=16, bold=True, color=KPMG_BLUE)
    desc = "Workbench 作为 \"AI 骨干系统\"，其核心价值在于解决全球 AI 开发的\"碎片化\"问题，通过统一平台实现三大战略价值：协作、一致性与规模化。"
    add_text_box(slide, 0.6, 2.3, 8.8, 0.7, desc, font_size=14, color=TEXT_PRIMARY)

    # 三大支柱
    cards_data = [
        {
            "title": "核心问题：碎片化",
            "desc": "解决全球网络中的重复劳动、成本激增和标准不一。",
            "left": 0.5
        },
        {
            "title": "核心定位：AI 骨干",
            "desc": "作为 \"AI Backbone\" 集中管理平台，赋能整个组织。",
            "left": 3.5
        },
        {
            "title": "核心价值：规模化",
            "desc": "推动全球协作、确保标准一致、实现规模化部署。",
            "left": 6.5
        }
    ]

    for card_data in cards_data:
        add_card(slide, card_data["left"], 3.3, 2.8, 2.5, bg_color=WHITE, border_color=BORDER_COLOR)
        add_text_box(slide, card_data["left"] + 0.15, 3.5, 2.5, 0.4, card_data["title"],
                     font_size=14, bold=True, color=TEXT_PRIMARY)
        add_text_box(slide, card_data["left"] + 0.15, 4.0, 2.5, 1.5, card_data["desc"],
                     font_size=12, color=TEXT_SECONDARY)

    # 页脚
    add_footer(slide, right_text="第 2 页 / 6")


def create_slide_3(prs):
    """幻灯片 3: 战略定位"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    add_text_box(slide, 0.5, 0.5, 9, 0.6, "战略定位：以 \"AI Backbone\" 解决碎片化问题",
                 font_size=36, bold=True, color=KPMG_BLUE)

    # 左侧：问题卡片
    add_card(slide, 0.5, 1.3, 4.5, 4.8, bg_color=RED_BG, border_color=RED_BORDER)
    add_text_box(slide, 0.65, 1.5, 4.2, 0.4, "问题：碎片化的 AI 开发 (Fragmentation)",
                 font_size=16, bold=True, color=RED_TEXT)
    add_text_box(slide, 0.65, 1.95, 4.2, 0.5, "各自为政的开发模式导致资源浪费和标准不一。",
                 font_size=12, color=TEXT_SECONDARY)

    problems = [
        "重复劳动 (Duplicated Work)\n不同成员公司重复发明相似的解决方案，缺乏共享。",
        "成本激增 (Increased Cost)\n无法形成规模效应，基础设施投资分散。",
        "标准不一致 (Inconsistent Standards)\nAI 应用质量参差不齐，缺乏统一的安全与合规标准。"
    ]

    y_pos = 2.6
    for problem in problems:
        add_text_box(slide, 0.65, y_pos, 4.2, 0.9, problem, font_size=11, color=TEXT_SECONDARY)
        y_pos += 1.1

    # 右侧：解决方案卡片
    add_card(slide, 5.2, 1.3, 4.5, 4.8, bg_color=GREEN_BG, border_color=GREEN_BORDER)
    add_text_box(slide, 5.35, 1.5, 4.2, 0.4, "解决方案：统一的 AI 骨干 (AI Backbone)",
                 font_size=16, bold=True, color=GREEN_TEXT)
    add_text_box(slide, 5.35, 1.95, 4.2, 0.5, "通过统一平台，实现 \"Build Once, Deploy Everywhere\"。",
                 font_size=12, color=TEXT_SECONDARY)

    solutions = [
        "全球协作 (Collaboration)\n全球资源驱动本地影响，任何创新成为全球共享资产。",
        "保证一致 (Consistency)\n统一的开发和 Trusted AI 标准，确保全球质量一致。",
        "实现规模 (Scale)\n快速复制成功方案，通过共享基础设施降低单位成本。"
    ]

    y_pos = 2.6
    for solution in solutions:
        add_text_box(slide, 5.35, y_pos, 4.2, 0.9, solution, font_size=11, color=TEXT_SECONDARY)
        y_pos += 1.1

    # 页脚
    add_footer(slide, right_text="第 3 页 / 6")


def create_slide_4(prs):
    """幻灯片 4: 技术架构"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    add_text_box(slide, 0.5, 0.4, 9, 0.5, "技术架构：三大差异化支柱",
                 font_size=36, bold=True, color=KPMG_BLUE)
    add_text_box(slide, 0.5, 0.95, 9, 0.4, "Workbench 通过三大技术支柱，建立区别于公有云 AI 的核心竞争优势。",
                 font_size=20, color=TEXT_SECONDARY)

    # 三个技术支柱卡片
    pillars = [
        {
            "title": "1. 全球数据主权",
            "desc": "承诺数据在指定地理区域内存储和处理，严格数据不出境。",
            "items": [
                "自动满足 GDPR 等各国法规",
                "保障敏感客户数据安全",
                "提供企业级数据主权承诺"
            ],
            "left": 0.5
        },
        {
            "title": "2. 可信 AI 框架",
            "desc": "\"Trusted AI Stamp\" 认证，将风控与合规能力嵌入 AI 全生命周期。",
            "items": [
                "管理 AI 特有风险（偏见、幻觉）",
                "适应全球 AI 监管法规演进",
                "建立客户信任的品牌标识"
            ],
            "left": 3.5
        },
        {
            "title": "3. 独特计费与遥测",
            "desc": "实现 API 级别的精准成本核算与使用量追踪。",
            "items": [
                "实现 \"按使用付费\"，成本透明化",
                "通过遥测数据进行使用模式分析",
                "建立内部市场机制，避免资源浪费"
            ],
            "left": 6.5
        }
    ]

    for pillar in pillars:
        add_card(slide, pillar["left"], 1.6, 2.8, 4.5, bg_color=WHITE, border_color=BORDER_COLOR)
        add_text_box(slide, pillar["left"] + 0.15, 1.8, 2.5, 0.3, pillar["title"],
                     font_size=14, bold=True, color=TEXT_PRIMARY)
        add_text_box(slide, pillar["left"] + 0.15, 2.2, 2.5, 0.6, pillar["desc"],
                     font_size=11, color=TEXT_SECONDARY)

        y_pos = 2.95
        for item in pillar["items"]:
            add_text_box(slide, pillar["left"] + 0.25, y_pos, 2.3, 0.4, "• " + item,
                         font_size=10, color=TEXT_SECONDARY)
            y_pos += 0.45

    # 页脚
    add_footer(slide, right_text="第 4 页 / 6")


def create_slide_5(prs):
    """幻灯片 5: 产品生态"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    add_text_box(slide, 0.5, 0.5, 9, 0.6, "产品生态：从平台能力到规模化应用",
                 font_size=36, bold=True, color=KPMG_BLUE)

    # 左侧：核心产品卡片
    add_card(slide, 0.5, 1.3, 6, 4.8, bg_color=WHITE, border_color=BORDER_COLOR)
    add_text_box(slide, 0.65, 1.5, 5.7, 0.4, "旗舰产品：aIQ Chat (KPMG 的企业版 ChatGPT)",
                 font_size=16, bold=True, color=KPMG_BLUE)
    desc = "在 KPMG 安全云环境中运行，可安全处理客户数据，并针对专业服务场景定制。"
    add_text_box(slide, 0.65, 1.95, 5.7, 0.5, desc, font_size=12, color=TEXT_SECONDARY)

    add_text_box(slide, 0.65, 2.6, 5.7, 0.3, "核心功能：知识增强 (Knowledge Grounding)",
                 font_size=13, bold=True, color=TEXT_PRIMARY)
    func_desc = "用户可上传团队文档或链接 SharePoint，创建定制化 AI 助手 (Personas)，确保 AI 响应基于内部知识库。"
    add_text_box(slide, 0.65, 2.95, 5.7, 0.8, func_desc, font_size=11, color=TEXT_SECONDARY)

    add_text_box(slide, 0.65, 3.85, 5.7, 0.3, "业务价值：多场景应用",
                 font_size=13, bold=True, color=TEXT_PRIMARY)
    value_desc = "赋能快速研究、政策解读、内容总结、文档起草，全面提高效率。"
    add_text_box(slide, 0.65, 4.2, 5.7, 0.6, value_desc, font_size=11, color=TEXT_SECONDARY)

    # 右侧：平台能力卡片
    add_card(slide, 6.7, 1.3, 3, 2.2, bg_color=LIGHT_GRAY, border_color=BORDER_COLOR)
    add_text_box(slide, 6.85, 1.5, 2.7, 0.3, "质量保证：Trusted AI Stamp",
                 font_size=13, bold=True, color=TEXT_PRIMARY)
    stamp_desc = "不仅是开发者工具，更是客户信任标志。将 KPMG 的审计和风控能力延伸至 AI 领域。"
    add_text_box(slide, 6.85, 1.85, 2.7, 1.4, stamp_desc, font_size=11, color=TEXT_SECONDARY)

    add_card(slide, 6.7, 3.7, 3, 2.2, bg_color=LIGHT_GRAY, border_color=BORDER_COLOR)
    add_text_box(slide, 6.85, 3.9, 2.7, 0.3, "平台化思维：赋能生态",
                 font_size=13, bold=True, color=TEXT_PRIMARY)
    platform_desc = "提供 AI 组件、能力 API 和统一设计系统，降低开发门槛，实现全球快速、标准化的应用部署。"
    add_text_box(slide, 6.85, 4.25, 2.7, 1.4, platform_desc, font_size=11, color=TEXT_SECONDARY)

    # 页脚
    add_footer(slide, right_text="第 5 页 / 6")


def create_slide_6(prs):
    """幻灯片 6: 总结页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    add_text_box(slide, 0.5, 0.4, 9, 0.5, "总结：Workbench 的核心战略价值",
                 font_size=36, bold=True, color=KPMG_BLUE)
    summary = "Workbench 是 KPMG 将其全球网络优势与专业服务基因深度融合的战略工具，是其 AI 时代的核心护城河。"
    add_text_box(slide, 0.5, 0.95, 9, 0.6, summary, font_size=18, color=TEXT_SECONDARY)

    # 左侧：三大战略意义
    add_card(slide, 0.5, 1.8, 4.5, 4.3, bg_color=WHITE, border_color=BORDER_COLOR)
    add_text_box(slide, 0.65, 2.0, 4.2, 0.4, "三大战略意义 (Strategic Implications)",
                 font_size=16, bold=True, color=KPMG_BLUE)

    implications = [
        "战略层面:\n从 \"碎片化\" 走向 \"AI 骨干\"，实现规模化创新。",
        "技术层面:\n建立 \"数据主权\" 和 \"可信 AI\" 的差异化护城河。",
        "应用层面:\n通过 aIQ Chat 等产品，将 AI 能力转化为实际业务价值。"
    ]

    y_pos = 2.6
    for impl in implications:
        add_text_box(slide, 0.75, y_pos, 4.0, 0.8, impl, font_size=12, color=TEXT_SECONDARY)
        y_pos += 1.1

    # 右侧：关键成功要素
    add_card(slide, 5.2, 1.8, 4.5, 4.3, bg_color=WHITE, border_color=BORDER_COLOR)
    add_text_box(slide, 5.35, 2.0, 4.2, 0.4, "关键成功要素 (Key Success Factors)",
                 font_size=16, bold=True, color=KPMG_BLUE)

    factors = [
        "全球网络优势 (实现网络效应)",
        "专业服务基因 (内嵌审计、风控能力)",
        "平衡创新与控制 (平台赋能 vs Trusted AI)",
        "技术与业务融合 (计费遥测系统支持运营)"
    ]

    y_pos = 2.6
    for factor in factors:
        add_text_box(slide, 5.45, y_pos, 4.0, 0.6, "• " + factor, font_size=12, color=TEXT_SECONDARY)
        y_pos += 0.8

    # 页脚
    add_footer(slide, right_text="第 6 页 / 6")


def main():
    """主函数"""
    print("开始生成 KPMG Workbench 深度概要 PowerPoint 文件...")

    # 创建演示文稿对象
    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9 宽屏
    prs.slide_height = Inches(7.5)

    # 创建 6 个幻灯片
    print("生成幻灯片 1: 封面页")
    create_slide_1(prs)

    print("生成幻灯片 2: 执行摘要")
    create_slide_2(prs)

    print("生成幻灯片 3: 战略定位")
    create_slide_3(prs)

    print("生成幻灯片 4: 技术架构")
    create_slide_4(prs)

    print("生成幻灯片 5: 产品生态")
    create_slide_5(prs)

    print("生成幻灯片 6: 总结页")
    create_slide_6(prs)

    # 确保输出目录存在
    output_dir = os.path.join(os.path.dirname(__file__), '..', 'generated_docs')
    os.makedirs(output_dir, exist_ok=True)

    # 保存文件
    output_path = os.path.join(output_dir, 'KPMG_Workbench深度概要.pptx')
    prs.save(output_path)

    print(f"\n✓ PowerPoint 文件已成功生成！")
    print(f"  输出路径: {os.path.abspath(output_path)}")
    print("\n提示: 请使用 Microsoft PowerPoint 或 WPS 演示打开文件查看效果。")


if __name__ == '__main__':
    main()
